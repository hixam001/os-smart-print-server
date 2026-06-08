import sys
import subprocess

try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "OS Smart Print Server Simulator"
subtitle.text = "Concurrency, Synchronization, and Scheduling Visualization"

# Slide 2: The Core OS Problem
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Core OS Problem"
tf = slide.placeholders[1].text_frame
tf.text = "A Networked Print Server perfectly simulates three OS challenges:"
p = tf.add_paragraph()
p.text = "1. Producer-Consumer Problem: Managing a bounded memory queue between fast users and slow printers."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Critical Sections: Preventing queue crashes when users submit jobs simultaneously."
p.level = 1
p = tf.add_paragraph()
p.text = "3. Job Scheduling: Deciding which job the printer should process next when the queue is full."
p.level = 1

# Slide 3: Architecture & Tech Stack
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Architecture & Tech Stack"
tf = slide.placeholders[1].text_frame
tf.text = "Decoupled 'Kernel Space' and 'User Space' Architecture:"
p = tf.add_paragraph()
p.text = "Backend (OS Kernel): Java 17 with Spring Boot. Uses native threading libraries to accurately simulate concurrency."
p.level = 1
p = tf.add_paragraph()
p.text = "Frontend (User Space): React.js and Vite. Provides a premium, 60fps 'glassmorphism' visualization."
p.level = 1
p = tf.add_paragraph()
p.text = "Communication: WebSockets act like hardware interrupts, sending state updates in sub-milliseconds."
p.level = 1

# Slide 4: Solving Concurrency: Semaphores
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Solving Concurrency: Semaphores"
tf = slide.placeholders[1].text_frame
tf.text = "Managing the Bounded Buffer:"
p = tf.add_paragraph()
p.text = "Custom Counting Semaphore manages available slots in the queue."
p.level = 1
p = tf.add_paragraph()
p.text = "Users (Producers) request permits. If full, threads suspend natively via Java's wait()."
p.level = 1
p = tf.add_paragraph()
p.text = "Printers (Consumers) release permits upon completion, triggering notifyAll() to wake suspended users."
p.level = 1

# Slide 5: Solving Concurrency: Locks
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Solving Concurrency: Mutex Locks"
tf = slide.placeholders[1].text_frame
tf.text = "Protecting Critical Sections:"
p = tf.add_paragraph()
p.text = "The Print Queue is a highly volatile shared data structure."
p.level = 1
p = tf.add_paragraph()
p.text = "A ReentrantLock (Mutex) surrounds all queue insertion and extraction logic."
p.level = 1
p = tf.add_paragraph()
p.text = "Guarantees absolute mutual exclusion, preventing race conditions, lost jobs, or duplicate reads."
p.level = 1

# Slide 6: CPU Scheduling Algorithms
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "CPU Scheduling Algorithms"
tf = slide.placeholders[1].text_frame
tf.text = "How the OS decides what prints next:"
p = tf.add_paragraph()
p.text = "FCFS: Fair, but suffers from the 'Convoy Effect' (small jobs stuck behind massive ones)."
p.level = 1
p = tf.add_paragraph()
p.text = "SJF: Fixes the Convoy Effect, but causes Absolute Starvation for large jobs."
p.level = 1
p = tf.add_paragraph()
p.text = "Hybrid (Aging): Uses SJF, but dynamically increases priority as a job waits. Guarantees no starvation."
p.level = 1

# Slide 7: Challenges Overcome
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Challenges Overcome"
tf = slide.placeholders[1].text_frame
tf.text = "Key technical hurdles during development:"
p = tf.add_paragraph()
p.text = "State Consistency: UI glitches under heavy load. Solved by forcing telemetry generation inside the Mutex lock for atomic snapshots."
p.level = 1
p = tf.add_paragraph()
p.text = "Starvation: SJF caused massive wait times. Solved mathematically using the Aging algorithm."
p.level = 1
p = tf.add_paragraph()
p.text = "Cross-Platform Deployment: Wrote custom scripts to bundle the full-stack app into a single executable .jar file."
p.level = 1

# Slide 8: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
tf = slide.placeholders[1].text_frame
tf.text = "Bridging the gap between theory and software engineering:"
p = tf.add_paragraph()
p.text = "Without mutual exclusion, a system collapses instantly under load."
p.level = 1
p = tf.add_paragraph()
p.text = "While simple algorithms like SJF look good, advanced concepts like Aging are strictly necessary."
p.level = 1
p = tf.add_paragraph()
p.text = "Abstract textbook concepts can be beautifully and effectively visualized in modern full-stack environments."
p.level = 1

prs.save("OS_Project_Presentation.pptx")
print("Presentation generated successfully!")
